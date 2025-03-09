import streamlit as st
from st_cookies_manager import EncryptedCookieManager
import csv
import re
from pdf2image import convert_from_bytes
import pytesseract
from PyPDF2 import PdfReader
from langchain_core.prompts import ChatPromptTemplate
from langchain_google_genai import GoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
import re
from duckduckgo_search import DDGS
from duckduckgo_search.exceptions import RatelimitException
from langchain_community.chat_models import ChatOpenAI
import base64
import requests
from google.api_core.exceptions import ResourceExhausted
import markdown
import pdfkit
from youtube_transcript_api import YouTubeTranscriptApi
import time
from pptx import Presentation
from docx import Document



def universal_setup(
    page_title="Home", page_icon="📝", upload_file_types=[], yt_upload=False, worker=False
):
    st.set_page_config(
        page_title=f"NoteCraft AI - {page_title}", page_icon=page_icon, layout="wide"
    )
    if page_title and page_title != "Home":
        st.header(f"NoteCraft AI - {page_title}")

    st.session_state["cookies"] = EncryptedCookieManager(
        prefix=st.secrets["COOKIES_PREFIX"],
        password=st.secrets["COOKIES_PASSWORD"],
    )
    if not st.session_state["cookies"].ready():
        st.stop()
    if worker and "worker" not in st.session_state:
        st.session_state["worker"] = LLMAgent(st.session_state["cookies"])
    if upload_file_types:
        if yt_upload:
            upload = st.sidebar.selectbox("Upload Type", ["File", "YouTube Video"])
            st.session_state["upload"] = ("file", st.sidebar.file_uploader(
                "upload your file", type=upload_file_types
            )) if upload == "File" else ("youtube", st.sidebar.text_input("YouTube Video URL"))
        else:
            st.session_state["file"] = st.sidebar.file_uploader(
                "upload your file", type=upload_file_types
            )



class LLMAgent:
    def __init__(self, cookies):
        self.cookies = cookies
        self.model = self.cookies["model"] if "model" in self.cookies else None
        self.llm = self._initialize_llm()

    def _initialize_llm(self):
        if self.model == "Gemini-1.5":
            llm = GoogleGenerativeAI(
                model="gemini-1.5-pro",
                api_key=self.cookies["GOOGLE_API_KEY"],
            )
        elif self.model == "GPT-4o-mini":
            llm = ChatOpenAI(
                model="gpt-4o-mini",
                api_key=self.cookies["OPENAI_API_KEY"],
            )
        else:
            st.write(
                f"Please make sure to select a model in the get access page, and you have set a valid API key."
            )
            st.stop()
        return llm

    def _get_chain(self, task):
        task_prompts = {
            "note_w_images": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        """You are a student writing notes from this transcript, Make sections headers, include all the main ideas in bullets and sub-bullets or in tables or images. Strictly base your notes on the provided information, without adding any external information. Try to make it clear and as simple as possible using simple terms and giving examples. Output in Markdown text formatting. To add images use this formatting: <<image search prompt>> use a general description that can be easily found in a Google search, avoid overly specific descriptions. Do it in {word_range} words.""",
                    ),
                    ("user", "{transcript}"),
                ]
            ),
            "note": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        """You are a student writing notes from this transcript, Make sections headers, include all the main ideas in bullets and sub-bullets or in tables. Strictly base your notes on the provided information, without adding any external information. Try to make it clear and as simple as possible using simple terms and giving examples. Output in Markdown text formatting. Do it in {word_range} words.""",
                    ),
                    ("user", "{transcript}"),
                ]
            ),
            "edit_note": ChatPromptTemplate.from_messages(
                [
                    ("system", """ you are tasked to edit this note: {text}."""),
                    (
                        "user",
                        "{request}\nOutput in Markdown formatting. to add images use this formatting: <<Write image search prompt here>> use a general description that can be easily found in a Google search, avoid overly specific descriptions.",
                    ),
                ]
            ),
            "page_note": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        "Write a brief summary of the following text, ensuring it does not exceed the word count of the original text and includes only the information present in the text—no additional details. here is the text: {transcript}",
                    ),
                ]
            ),
            "edit_flashcards": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        """ You are tasked to make an edit to these flashcards: {text}.""",
                    ),
                    (
                        "user",
                        "{request}\nOutput in the same csv formate with '\t' as the seperator. each row should contain 2 columns: Question \t Answer. only return the csv data without any other information.",
                    ),
                ]
            ),
            "Term --> Definition": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        "You are tasked with creating flashcards that will help students learn the important terms, proper nouns and concepts in this note. Only make flashcards directly related to the main idea of the note, include as much detail as possible in each flashcard, returning it in a CSV formate with '\t' as the seperator. each row should include 2 columns: Term \t Definition. make exactly from {flashcard_range} flashcards. only return the csv data without any other information.",
                    ),
                    ("user", "{transcript}"),
                ]
            ),
            "Question --> Answer": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        "You are tasked with creating a mock test that will help students learn and understand concepts in this note. Only make questions directly related to the main idea of the note, You should include all these question types: fill in the blank, essay questions and short answer questions.. return the questions and answers in a CSV formate with '\t' as the seperator. Each row should include 2 columns: question \t answer. make exactly from {flashcard_range} Questions. only return the csv data without any other information.",
                    ),
                    ("user", "{transcript}"),
                ]
            ),
            "MCQ --> Answer": ChatPromptTemplate.from_messages(
                [
                    (
                        "system",
                        "You are tasked with creating a mock test with multiple choice questions that will help students learn and understand concepts in this note. Only make questions directly related to the main idea of the note. Include multiple choice and True/False Questions. return the questions and answers in a CSV formate with '\t' as the seperator. Each row should include 2 columns: question \t Answer. To seperate the question from the options in use <br>. make exactly from {flashcard_range} Questions. only return the csv data without any other information.",
                    ),
                    ("user", "{transcript}"),
                ]
            ),
        }
        prompt = task_prompts.get(task)
        if prompt is None:
            raise ValueError(f"Task '{task}' not found in task_prompts.")
        return prompt | self.llm

    def get_note(self, transcript, word_range: tuple, images=False):
        note_prompt = self._get_chain("note" if not images else "note_w_images")
        if st.session_state["cookies"]["NoteForge"] == "True":
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=3000,
                chunk_overlap=200,
                separators=["\n", ".", "!", "?"],
            )
            transcript = text_splitter.split_text(transcript)
            page_notes_chain = self._get_chain("page_note")
            try:
                final_transcript = "\n".join(
                    (
                        page_notes_chain.invoke({"transcript": doc})
                        if st.session_state["cookies"]["model"] == "Gemini-1.5"
                        else page_notes_chain.invoke({"transcript": doc}).content
                    )
                    for doc in transcript
                )
            except ResourceExhausted:
                st.error(
                    "API Exhausted, if you are using the free version of the API, you may have reached the limit.\nTry again later.\nIf NoteForge is enabled, try disabling it."
                )
                st.stop()
        else:
            final_transcript = transcript
        try:

            self.note = note_prompt.invoke(
                {
                    "transcript": final_transcript,
                    "word_range": " to ".join(map(str, word_range)),
                }
            )
        except ResourceExhausted:
            st.error(
                "API Exhausted, if you are using the free version of the API, you may have reached the limit.\nTry again later.\nIf NoteForge is enabled, try disabling it."
            )
            st.stop()
        return str(self.note) if st.session_state["cookies"]["model"] == "Gemini-1.5" else str(self.note.content)

    def get_flashcards(
        self, flashcard_range: tuple, task="Term --> Definition", transcript=None
    ):
        if (
            st.session_state["cookies"]["NoteForge"] == "True"
            and transcript is not None
        ):
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=3000,
                chunk_overlap=200,
                separators=["\n", ".", "!", "?"],
            )
            transcript = text_splitter.split_text(transcript)
            page_notes_chain = self._get_chain("page_note")
            flashcard_chain = self._get_chain(task)
            try:
                final_transcript = "\n".join(
                    (
                        page_notes_chain.invoke({"transcript": doc})
                        if st.session_state["cookies"]["model"] == "Gemini-1.5"
                        else page_notes_chain.invoke({"transcript": doc}).content
                    )
                    for doc in transcript
                )
            except ResourceExhausted:
                st.error(
                    "API Exhausted, if you are using the free version of the API, you may have reached the limit.\nTry again later.\nIf NoteForge is enabled, try disabling it."
                )
                st.stop()
        elif transcript is None:
            final_transcript = self.note
        elif transcript is not None:
            final_transcript = transcript
        try:
            flashcard_chain = self._get_chain(task)
            self.flashcards = flashcard_chain.invoke(
                {
                    "transcript": final_transcript,
                    "flashcard_range": " to ".join(map(str, flashcard_range)),
                }
            )
        except ResourceExhausted:
            st.error(
                "API Exhausted, if you are using the free version of the API, you may have reached the limit.\nTry again later.\nIf NoteForge is enabled, try disabling it."
            )
            st.stop()
        return str(self.flashcards) if st.session_state["cookies"]["model"] == "Gemini-1.5" else str(self.flashcards.content)

    def edit(self, task, request, text):
        edit_chain = self._get_chain(task)
        try:
            return str(edit_chain.invoke({"request": request, "text": text}).content)
        except ResourceExhausted:
            st.error(
                "API Exhausted, if you are using the free version of the API, you may have reached the limit.\nTry again later.\nIf NoteForge is enabled, try disabling it."
            )
            st.stop()


from duckduckgo_search.exceptions import RatelimitException

def md_image_format(md, encoded=False):
    def replace_with_image(match):
        description = match.group(1).strip()
        retry_attempts = 5
        for attempt in range(retry_attempts):
            try:
                results = ddgs.images(keywords=description, max_results=5)
                for result in results:
                    image_url = result["image"]
                    if encoded:
                        try:
                            response = requests.get(image_url)
                            response.raise_for_status()
                            image_data = response.content
                            image_format = image_url.split(".")[-1]
                            base64_image = base64.b64encode(image_data).decode("utf-8")
                            return f"![{description}](data:image/{image_format};base64,{base64_image})"
                        except requests.RequestException as e:
                            continue
                    else:
                        return f"![{description}]({image_url})"
                return "\n"
            except RatelimitException:
                if attempt < retry_attempts - 1:
                    time.sleep(2 ** attempt)  # Exponential backoff
                else:
                    st.error("Rate limit exceeded. Please try again later.")
                    st.stop()

    pattern = r"<<\s*(.*?)\s*>>"
    ddgs = DDGS()
    modified_md = re.sub(pattern, replace_with_image, str(md), flags=re.DOTALL)
    return modified_md


def get_base64_encoded_pdf(file):
    file.seek(0)
    pdf_content = file.read()
    encoded_pdf = base64.b64encode(pdf_content).decode("utf-8")
    return encoded_pdf


def get_document_text(file, page_range: tuple = None):
    file_extension = file.name.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        try:
            pdf_reader = PdfReader(file)
        except Exception as e:
            return f"There was a problem reading the pdf: {str(e)}"

        if page_range is None:
            first_page = 1
            last_page = len(pdf_reader.pages)
        else:
            first_page, last_page = page_range

        text = ""
        for page_num in range(first_page - 1, last_page):
            page = pdf_reader.pages[page_num]
            if "/XObject" in page["/Resources"]:
                xObject = page["/Resources"]["/XObject"].get_object()
                if any(xObject[obj]["/Subtype"] == "/Image" for obj in xObject):
                    images = convert_from_bytes(
                        file.getvalue(), first_page=page_num + 1, last_page=page_num + 1
                    )
                    for image in images:
                        image_text = pytesseract.image_to_string(image)
                        if image_text.strip():
                            text += image_text + "\n"
                else:
                    text += page.extract_text() + "\n"
            else:
                text += page.extract_text() + "\n"
        return text

    elif file_extension == 'pptx':
        try:
            presentation = Presentation(file)
        except Exception as e:
            return f"There was a problem reading the pptx: {str(e)}"

        if page_range is None:
            first_slide = 1
            last_slide = len(presentation.slides)
        else:
            first_slide, last_slide = page_range

        text = ""
        for slide_num in range(first_slide - 1, last_slide):
            slide = presentation.slides[slide_num]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif file_extension == 'docx':
        try:
            document = Document(file)
        except Exception as e:
            st.error(f"There was a problem reading the docx: {str(e)}")
            st.stop()
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        return text

    else:
        st.write("Unsupported file type")
        st.stop()

def fetch_transcript(url):
    if "watch?v=" in url:
        video_id = url.split("watch?v=")[1].split("&")[0]
    elif "youtu.be/" in url:
        video_id = url.split("youtu.be/")[1].split("?")[0]
    else:
        st.error("Invalid YouTube URL")
        st.stop()
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        return "\n".join([item['text'] for item in transcript])
    except:
        st.error(f"Error fetching transcript, the video might have disabled captions.")
        st.stop()

def page_count(document):
    """Extract the number of pages in a document (PDF, DOCX, PPTX)."""
    file_extension = document.name.split('.')[-1].lower()
    
    try:
        if file_extension == 'pdf':
            pdf_reader = PdfReader(document)
            return len(pdf_reader.pages)
        elif file_extension == 'docx':
            return 1
        elif file_extension == 'pptx':
            presentation = Presentation(document)
            return len(presentation.slides)
        else:
            st.error("Unsupported file type")
            st.stop()
    except Exception as e:
        st.error(f"Error reading the {file_extension.upper()} file: {str(e)}")
        st.stop()


def clean_flashcards(flashcards):
    unwanted_headers_col1 = {"question", "questions", "term", "terms"}
    unwanted_headers_col2 = {"answer", "answers", "definition", "definitions"}

    flashcards = re.sub(r"^\s*```(?:csv\s*)?", "", flashcards, flags=re.IGNORECASE)
    flashcards = re.sub(r"```\s*$", "", flashcards, flags=re.IGNORECASE)

    pattern = re.compile(r"^\s*(\w+)\s*\t\s*(\w+)\s*\n", re.IGNORECASE)
    match = pattern.match(flashcards)

    if match:
        col1, col2 = match.groups()
        if (
            col1.lower() in unwanted_headers_col1
            and col2.lower() in unwanted_headers_col2
        ):
            flashcards = flashcards[match.end() :]

    flashcards = re.sub(r"\t+", "\t", flashcards)

    return flashcards.strip()


def display_flashcards(flashcards):
    if not flashcards:
        st.error("Error Generating Flashcards, No Flashcards generated.")
        return

    flashcards = clean_flashcards(flashcards)
    try:
        flashcards_reader = csv.reader(flashcards.splitlines(), delimiter="\t")
        flashcards_data = list(flashcards_reader)
    except csv.Error:
        st.error("There was an error generating the flashcards, please try again.")
        st.stop()

    st.session_state.questions = [row[0] for row in flashcards_data if len(row) > 1]
    st.session_state.answers = [row[1] for row in flashcards_data if len(row) > 1]

    if "current_question_index" not in st.session_state:
        st.session_state.current_question_index = 0
    if "show_answer" not in st.session_state:
        st.session_state.show_answer = False

    col1, col2, col3 = st.columns(3, gap="small")
    if col1.button("Previous Question"):
        if st.session_state.current_question_index > 0:
            st.session_state.current_question_index -= 1
            st.session_state.show_answer = False
        else:
            st.warning("You are at the first question.")

    if col2.button("Show Answer"):
        st.session_state.show_answer = True

    st.write(f"Total Flashcards {len(st.session_state.questions)}.")

    if col3.button("Next Question"):
        if (
            st.session_state.current_question_index
            < len(st.session_state.questions) - 1
        ):
            st.session_state.current_question_index += 1
            st.session_state.show_answer = False
        else:
            st.warning("You have reached the last question.")

    question = st.session_state.questions[st.session_state.current_question_index]
    st.html(f"Question {st.session_state.current_question_index + 1}: {question}")

    if st.session_state.show_answer:
        answer = st.session_state.answers[st.session_state.current_question_index]
        st.write(f"Answer: {answer}")


def parse_studkit(content):
    note = re.search(r"note=`\^(.*?)`\^", content, re.DOTALL).group(1)
    flashcards = re.search(r"flashcards=`\^(.*?)`\^", content, re.DOTALL).group(1)
    return note, flashcards


def paper(header_text, markdown_text=None, flashcards=None):
    styles = """
body { font-family: serif; }
h1 { text-align: center; font-size: 36px; }
.page-headers { text-align: center; font-size: 24px; margin-bottom: 25px; }
a { color: black; }
.cover_page_title { text-align: center; margin-top: 30%; font-size: 70px; }
h1, h2, h3, h4, h5, h6 { margin-top: 50px; }
p, li { font-size: 20px; }
li br { margin-bottom: 20px; display: block; content: "";}
table { border-collapse: collapse; width: 100%; margin: 1em 0; }
th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
th { background-color: #f5f5f5; }
blockquote { border-left: 4px solid #ff4b4b; margin: 0; padding-left: 20px; color: #666; }
img { max-width: 90%; height: auto; margin: 20px auto; display: block; }
pre { background: #cfd6e3; padding: 10px; border: 1px solid #ddd; border-radius: 5px; margin: 0 40px; }
code { font-family: 'Courier New', Courier, monospace; }
.questions li { padding-bottom: 150px; }
.answers li { margin-bottom: 15px; }
"""
    cover_html = f"""
    <html>
    <head>
        <title>Cover Page</title>
        <style>
            {styles}
        </style>
    </head>
    <body>
        <h1 class="cover_page_title">{header_text}</h1>
        <p style="text-align: center; font-size: 24px;">Crafted by <a href="https://notecraft.streamlit.app/">NoteCraft</a></p>
        <p style="text-align: center; font-size: 24px;"><svg xmlns="http://www.w3.org/2000/svg" width="25" height="25" fill="currentColor" class="bi bi-github" viewBox="0 0 16 16">
  <path d="M8 0C3.58 0 0 3.58 0 8c0 3.54 2.29 6.53 5.47 7.59.4.07.55-.17.55-.38 0-.19-.01-.82-.01-1.49-2.01.37-2.53-.49-2.69-.94-.09-.23-.48-.94-.82-1.13-.28-.15-.68-.52-.01-.53.63-.01 1.08.58 1.23.82.72 1.21 1.87.87 2.33.66.07-.52.28-.87.51-1.07-1.78-.2-3.64-.89-3.64-3.95 0-.87.31-1.59.82-2.15-.08-.2-.36-1.02.08-2.12 0 0 .67-.21 2.2.82.64-.18 1.32-.27 2-.27s1.36.09 2 .27c1.53-1.04 2.2-.82 2.2-.82.44 1.1.16 1.92.08 2.12.51.56.82 1.27.82 2.15 0 3.07-1.87 3.75-3.65 3.95.29.25.54.73.54 1.48 0 1.07-.01 1.93-.01 2.2 0 .21.15.46.55.38A8.01 8.01 0 0 0 16 8c0-4.42-3.58-8-8-8"/></svg> GitHub repo: https://github.com/AL-Sayed1/NoteCraft</p>
    </body>
    </html>
    """

    html_text = ""
    if markdown_text:
        html_text = (
            "<div style='page-break-after: always;'></div> <h1 class='.page-headers'>Notes:</h1>"
            + markdown.markdown(
                markdown_text,
                extensions=[
                    "tables",
                    "fenced_code",
                    "attr_list",
                    "toc",
                    "footnotes",
                    "codehilite",
                    "meta",
                ],
            )
        )

    questions_html = ""
    answers_html = ""
    if flashcards:
        questions_html = f"""
        <div style='page-break-after: always;'></div>
        <html>
        <head>
            <title>Questions</title>
            <style>
                {styles}
            </style>
        </head>
        <body>
            <h1 class='page-headers'>Questions</h1>
            <ol class="questions">
        """
        answers_html = f"""
        <div style='page-break-after: always;'></div>
        <html>
        <head>
            <title>Answer Key</title>
            <style>
                {styles}
            </style>
        </head>
        <body>
            <h1 class='page-headers'>Answer Key</h1>
            <ol class="answers">
        """

        flashcards_reader = csv.reader(
            clean_flashcards(flashcards).splitlines(), delimiter="\t"
        )
        for row in flashcards_reader:
            if len(row) == 2:
                question, answer = row
                questions_html += f"<li>{question}</li>"
                answers_html += f"<li>{answer}</li>"
            else:
                questions_html += f"<li>Invalid question format: {' '.join(row)}</li>"
                answers_html += f"<li>Invalid answer format: {' '.join(row)}</li>"

        questions_html += "</ol></body></html>"
        answers_html += "</ol></body></html>"

    full_html = cover_html + html_text + questions_html + answers_html
    pdf_data = pdfkit.from_string(
        full_html,
        False,
        options={
            "page-size": "A4",
            "margin-top": "2cm",
            "margin-right": "2cm",
            "margin-bottom": "2cm",
            "margin-left": "2cm",
            "encoding": "UTF-8",
            "no-outline": None,
        },
    )
    return pdf_data
